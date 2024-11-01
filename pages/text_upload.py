import streamlit as st

import io
import apis.aws_api as aws_api
import apis.dify_api as dify_api

import mimetypes
import fitz
import pandas as pd
from docx import Document
from pptx import Presentation

from botocore.exceptions import NoCredentialsError

def patched_close():
    print("good try boto3")

def preview_txt(file):
    content = file.read().decode("utf-8")
    st.text_area("File Preview", content, height=300, disabled=True)

def preview_pdf(file):
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")    
    # Display the first page
    page = pdf_document.load_page(0)  # Load the first page (index 0)
    pix = page.get_pixmap()  # Render page to image
    img = io.BytesIO(pix.tobytes())  # Convert to BytesIO
    st.image(img, caption="Page 1", use_column_width=True)

def preview_docx(file):
    doc = Document(file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    complete_text = "\n".join(full_text)
    # Limit the text to a certain number of characters (e.g., 500)
    limited_text = complete_text[:500] + ("..." if len(complete_text) > 500 else "")
    # Display the limited content in a text_area
    st.text_area("File Preview", value=limited_text, height=300, disabled=True)

def preview_csv(file):
    df = pd.read_csv(file)
    st.write("### CSV Content")
    st.dataframe(df, height=300)

def preview_xlsx(file):
    df = pd.read_excel(file)
    st.write("### XLSX Content")
    st.dataframe(df, height=300)

def preview_pptx(file):
    prs = Presentation(uploaded_file)
    
    first_slide = prs.slides[0]
    
    content = ""
    for shape in first_slide.shapes:
        if hasattr(shape, "text") and shape.text:
            content += shape.text + '\n'
    st.text_area("File Preview", value=content, height=300, disabled=True)

# Streamlit UI
st.title("Text Upload to AWS S3 Folder")

# File upload widget
uploaded_file = st.file_uploader("Choose an text file", type=["txt", "pdf", "docx", "csv", "xlsx", "pptx"])

if uploaded_file is not None:
    # Get file format and name
    mime_type, _ = mimetypes.guess_type(uploaded_file.name)
    print(mime_type)
    if mime_type == 'text/plain':
        preview_txt(uploaded_file)
    elif mime_type == "application/pdf":
        preview_pdf(uploaded_file)
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        preview_docx(uploaded_file)
    elif mime_type == "application/vnd.ms-excel":
        preview_csv(uploaded_file)
    elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        preview_xlsx(uploaded_file)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        preview_pptx(uploaded_file)
    uploaded_file.seek(0)

    # Upload text file to S3 when button is clicked
    if st.button('Upload to S3'):
        # Upload the image as binary stream (BytesIO object) to S3        
        try:
            text_url = aws_api.upload_text_to_s3(uploaded_file, mime_type, aws_api.TEXT_BUCKET_NAME, aws_api.NEWS_FOLDER_NAME)
            st.success(f"File uploaded successfully to {aws_api.TEXT_BUCKET_NAME}")
            
            dify_code, dify_msg = dify_api.text_upload_subprocess(uploaded_file, mime_type)
            
            if dify_code == 0:
                st.write("Text information synchronized successfully in Dify")
            else:
                st.write(dify_msg)
        except NoCredentialsError:
            st.error("Credentials not available")

if st.button("Back to Home"):
    st.switch_page("homepage.py")